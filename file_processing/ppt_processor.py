from pptx import Presentation


class PPTProcessor:

    @staticmethod
    def extract_text(ppt_path):
        """
        Extract text from PPTX file
        """

        text = ""

        try:
            presentation = Presentation(ppt_path)

            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

            return text

        except Exception as e:
            print(f"Error: {e}")
            return ""